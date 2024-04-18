#!/usr/bin/python3
# -*- coding: utf-8 -*-
#
# Based on Bishwajyotichaudhary's Meduium blog "Gemini Presentation Maker Simple Python Implementation" on Jan 3, 2024
#
# See his github: https://github.com/bishwajyotichaudhary46/AI-Presentation-Maker-Gemini-Pro
#
# run with: streamlit run new_app.py
#
# Extentions inspired by material from Sam Witteveen's video "Clause 3 Vs Gemini Vs GPT-4: What Can Make Amazing Powerpoints?": https://www.youtube.com/watch?v=g3vroajaWwg
#
# When using the KTH template - besure that you have isntalled the Figtree font
# see https://fonts.google.com/specimen/Figtree
#
# KTH template - layouts:
# 0 Title, KTH blue
# 1 Title, navy blue
# 2 Title, image
# 3 Title/chapter, light blue
# 4 Title/chapter, image, light blue
# 5 Title/chapter, sand
# 6 Title/chapter, image, sand
# 7 Title/chapter, KTH blue
# 8 Title/chapter, image, KTH blue
# 9 Title/chapter, navy blue
# 10 Title/chapter, image, navy blue
# 11 Headline and content
# 12 Headline and content, blue lines
# 13 Headline and content, sand lines
# 14 Headline and content, light blue 
# 15 Headline and content, sand
# 16 Two parts
# 17 Two parts, subheader
# 18 Text and image
# 19 Comparison, light blue
# 20 Comparison, sand
# 21 Comparison, KTH blue
# 22 Comparison, navy blue
# 23 Comparison, green
# 24 Comparison, teal
# 25 Comparison, red
# 26 Comparison, yellow
# 27 Headline and 1 image
# 28 Headline and 2 images
# 29 Headline and 2 images, text
# 30 Headline and 3 images
# 31 Headline and 3 images, text
# 32 Headline and 4 images
# 33 Headline and 4 images, text
# 34 Headline only
# 35 Headline only, light blue lines
# 36 Headline only, sand lines
# 37 Headline only, light blue
# 38 Headline only, sand
# 39 Empty
# 40 Headline and content, top
# 41 Headline only, top
# 42 Logo
# 43 Logo, navy blue
# 44 How to use template
# 45 Accessibility Office 2013/2016
# 46 Accessibility Office 365/2019

from dotenv import load_dotenv
load_dotenv()

import os
import streamlit as st
import google.generativeai as genai
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re

from pptx import Presentation


os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_gemini_response(prompt):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(prompt)
    return response

sub_titles = []
def refine_subtopics(sub_topics, sub_titles):
    for sub_topic in sub_topics:
        sub_titles.append(sub_topic[3:].replace('"',""))
    return sub_titles

content = []
def content_generation(sub_titles):
    for i in sub_titles:
        prompt = f"Generate a content of {i} for presentation slide on the 2 bullet point only each of point 20 tokens"
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content(prompt)
        content.append(response.text)
    return content

def clean_text(text):
    # Remove extra whitespaces and newlines
    cleaned_text = re.sub('\s+', ' ', text).strip()

    # Remove markdown-style bullet points, asterisks, and numeric bullet points
    cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)

    # Remove extra spaces before and after colons
    cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)

    # Remove extra spaces before and after hyphens
    cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)

    return cleaned_text


def split_sentences(text):
    # Split the text into sentences using regular expression
    sentences = re.split(r'(?<=\.)\s+', text)

    # Capitalize the first letter of each sentence
    sentences = [sentence.capitalize() for sentence in sentences]

    return sentences

def replace_and_capitalize(text):
    # Define a function to replace and capitalize the text between colons
    def replace_and_capitalize_colon(match):
        return match.group(1) + match.group(2).capitalize() + match.group(3)

    # Use regular expression to find and replace text between colons
    result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)', replace_and_capitalize_colon, text)

    return result


final_content = []
def refine_final_content(content):
    for i in content:
        cleaned_text = clean_text(i)
        sentences = split_sentences(cleaned_text)
        final_content.append(sentences)
    print("final content ready....")
    return final_content

powerpoint = Presentation("Presentation1.pptx") # use the KTH template

def slide_maker(powerpoint, topic,sub_titles, final_content):
    title_slide_layout = powerpoint.slide_layouts[0]
    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]
    content.text = "Created By AI Gemini Model"

    # Get a list of the slide layouts that are available and their index
    for idx, layout in enumerate(powerpoint.slide_layouts):
        print(idx, layout.name)

    for i in range(len(sub_titles)):
        bulletLayout = powerpoint.slide_layouts[11]
        secondSlide = powerpoint.slides.add_slide(bulletLayout)
        # accessing the attributes of shapes
        myShapes = secondSlide.shapes
        titleShape = myShapes.title
        bodyShape = myShapes.placeholders[1]
        titleShape.text = sub_titles[i]
        titleShape.text_frame.paragraphs[0].font.size = Pt(24)
        titleShape.text_frame.paragraphs[0].font.bold = True
        tFrame = bodyShape.text_frame
        print("Topic Generated")
        for point in final_content[i]:
            point = re.sub(r':[^:]+:', ':', point)
            point = replace_and_capitalize(point)
            p = tFrame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.space_after = Pt(10)
    return powerpoint
def download_button(file_path,topic):
    # Read the content of the PPT file
    with open(file_path, "rb") as file:
        ppt_content = file.read()

    # Create a download button for the PPT file
    st.download_button(
        label="Download PowerPoint",
        data=ppt_content,
        file_name=f"{topic}.pptx",  # Change the file name as needed
        key="ppt_download_button"
    )

#model = genai.GenerativeModel('gemini-pro')
st.set_page_config(page_title="Gemini Presentation Maker")

st.header("Gemini Presentation Maker")
topic=st.text_input("Input Prompt: ",key="input")
no_of_slide=st.text_input("Enter Number Of Slide: ",key="slide")

submit=st.button("Generate")
if submit:
    prompt =f"Generate a {no_of_slide} sub-titles only  on the topic of {topic}"
    response = get_gemini_response(prompt)
    print("Topic Generated")
    sub_topics = response.text.split("\n")
    sub_titles = refine_subtopics(sub_topics, sub_titles)
    print("Sub Titles")
    content = content_generation(sub_titles)
    print("content Generated")
    final_content = refine_final_content(content)
    #cleaned_text = clean_text(content[0])
    #sentences = split_sentences(cleaned_text)
    print("final content ready")
    powerpoint = slide_maker(powerpoint,topic, sub_titles, final_content)
    powerpoint.save(f"../Powerpoint/{topic}.pptx")
    st.text("Presentation Ready")
    download_button(f"../Powerpoint/{topic}.pptx",topic)
    print("Presentation Ready")
